"""Day 3. PowerPoint export.

Produces a short board-pack deck from an AnalysisResult:

* Cover slide: company, period, run timestamp.
* Headline KPIs: budget, actual, variance, variance %, RAG counts.
* AI commentary: headline + summary + adverse / favourable drivers + actions.
* Per-cost-centre slides (capped at 8 to keep the deck readable).

Skip silently when python-pptx is not installed; the rest of the app
keeps working. The Excel pack is the source of truth, the deck is a
nice-to-have.
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


INK    = (0x14, 0x21, 0x3A)
PAPER  = (0xFB, 0xF6, 0xEC)
BURG   = (0x8E, 0x2A, 0x2A)
AMBER  = (0xC1, 0x8B, 0x2A)
OLIVE  = (0x5E, 0x6E, 0x3B)
STEEL  = (0x1F, 0x4E, 0x79)
MUTED  = (0x6E, 0x66, 0x5A)

MAX_COST_CENTRE_SLIDES = 8


def is_available() -> bool:
    return _PPTX_AVAILABLE


def write_pptx(result, out_dir: Path | str) -> Path | None:
    """Render the deck. Returns the output path, or None if python-pptx is missing."""
    if not _PPTX_AVAILABLE:
        return None
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    company = _slug(result.metadata.company)
    period = _slug(result.metadata.period_label)
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    out = out_dir / f"budget_variance_{company}_{period}_{ts}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover(prs, result)
    _add_headline(prs, result)
    if not result.commentary.skipped and not result.commentary.error:
        _add_commentary(prs, result)
    _add_cost_centre_slides(prs, result)
    _add_outro(prs, result)

    prs.save(out)
    return out


# ---- Slide helpers ----------------------------------------------------

def _add_cover(prs, result) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, INK)
    _add_text(
        slide,
        result.metadata.company,
        left=Inches(0.7), top=Inches(2.6), width=Inches(12), height=Inches(1.5),
        size=44, bold=True, colour=PAPER, family="Cambria",
    )
    _add_text(
        slide,
        f"Budget vs actual. {result.metadata.period_label}",
        left=Inches(0.7), top=Inches(4.0), width=Inches(12), height=Inches(0.6),
        size=22, colour=PAPER, family="Cambria",
    )
    _add_text(
        slide,
        f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  "
        f"Currency {result.metadata.currency}  |  "
        f"Fiscal year {result.metadata.fiscal_year}",
        left=Inches(0.7), top=Inches(4.7), width=Inches(12), height=Inches(0.4),
        size=12, colour=(0xC0, 0xB8, 0xA0), italic=True,
    )


def _add_headline(prs, result) -> None:
    h = result.headline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, "Headline", Inches(0.6), Inches(0.4), Inches(8), Inches(0.6),
              size=28, bold=True, colour=INK, family="Cambria")

    rows = [
        ("Total budget",    _gbp(h.total_budget),    INK),
        ("Total actual",    _gbp(h.total_actual),    INK),
        ("Total variance",  _gbp_signed(h.total_variance),
                            BURG if h.total_variance > 0 else (STEEL if h.total_variance < 0 else INK)),
        ("Variance %",      _pct_signed(h.total_variance_pct), INK),
        ("Rows",            f"{h.row_count}",        MUTED),
        ("Periods",         f"{h.period_count}",     MUTED),
    ]
    for i, (label, value, colour) in enumerate(rows):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(1.4 + row * 2.2)
        _add_text(slide, label, x, y, Inches(4), Inches(0.4),
                  size=11, colour=MUTED, family="Cascadia Code")
        _add_text(slide, value, x, y + Inches(0.4), Inches(4), Inches(1.4),
                  size=36, bold=True, colour=colour, family="Cambria")

    # RAG counts strip along the bottom.
    bands = [
        ("RED",        h.rag_counts.get("red", 0),        BURG),
        ("AMBER",      h.rag_counts.get("amber", 0),      AMBER),
        ("GREEN",      h.rag_counts.get("green", 0),      OLIVE),
        ("FAVOURABLE", h.rag_counts.get("favourable", 0), STEEL),
        ("N/A",        h.rag_counts.get("na", 0),         MUTED),
    ]
    width = 2.4
    for i, (label, count, colour) in enumerate(bands):
        x = Inches(0.6 + i * width)
        y = Inches(6.2)
        _add_text(slide, str(count), x, y, Inches(width), Inches(0.6),
                  size=24, bold=True, colour=colour, family="Cambria")
        _add_text(slide, label, x, y + Inches(0.55), Inches(width), Inches(0.4),
                  size=10, colour=MUTED, family="Cascadia Code")


def _add_commentary(prs, result) -> None:
    c = result.commentary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, "Management commentary", Inches(0.6), Inches(0.4),
              Inches(12), Inches(0.6), size=28, bold=True, colour=INK, family="Cambria")
    _add_text(slide, c.headline, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
              size=20, bold=True, colour=INK, family="Cambria")
    _add_text(slide, c.summary, Inches(0.6), Inches(2.2), Inches(12), Inches(1.6),
              size=14, colour=INK)

    _add_text(slide, "Adverse drivers", Inches(0.6), Inches(4.0),
              Inches(6), Inches(0.4), size=12, bold=True, colour=BURG, family="Cascadia Code")
    _add_text(slide, _bullets(c.adverse_drivers), Inches(0.6), Inches(4.4),
              Inches(6), Inches(2.4), size=12, colour=INK)

    _add_text(slide, "Favourable drivers", Inches(7.0), Inches(4.0),
              Inches(6), Inches(0.4), size=12, bold=True, colour=STEEL, family="Cascadia Code")
    _add_text(slide, _bullets(c.favourable_drivers), Inches(7.0), Inches(4.4),
              Inches(6), Inches(2.4), size=12, colour=INK)


def _add_cost_centre_slides(prs, result) -> None:
    aggs = list(result.by_cost_centre)[:MAX_COST_CENTRE_SLIDES]
    for a in aggs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(slide, a.key, Inches(0.6), Inches(0.4),
                  Inches(12), Inches(0.6), size=28, bold=True, colour=INK, family="Cambria")
        _add_text(slide, "Cost centre", Inches(0.6), Inches(1.0),
                  Inches(12), Inches(0.4), size=11, colour=MUTED, family="Cascadia Code")

        kpis = [
            ("Budget",    _gbp(a.budget),                 INK),
            ("Actual",    _gbp(a.actual),                 INK),
            ("Variance",  _gbp_signed(a.variance),
                          BURG if a.variance > 0 else (STEEL if a.variance < 0 else INK)),
            ("Variance %",_pct_signed(a.variance_pct),   INK),
        ]
        for i, (label, value, colour) in enumerate(kpis):
            x = Inches(0.6 + i * 3.1)
            y = Inches(1.8)
            _add_text(slide, label, x, y, Inches(3), Inches(0.4),
                      size=11, colour=MUTED, family="Cascadia Code")
            _add_text(slide, value, x, y + Inches(0.4), Inches(3), Inches(1.4),
                      size=30, bold=True, colour=colour, family="Cambria")

        # RAG breakdown strip
        bands = [
            ("Red",         a.rag_counts.get("red", 0),        BURG),
            ("Amber",       a.rag_counts.get("amber", 0),      AMBER),
            ("Green",       a.rag_counts.get("green", 0),      OLIVE),
            ("Favourable",  a.rag_counts.get("favourable", 0), STEEL),
            ("n/a",         a.rag_counts.get("na", 0),         MUTED),
        ]
        for i, (label, count, colour) in enumerate(bands):
            x = Inches(0.6 + i * 2.4)
            y = Inches(4.4)
            _add_text(slide, str(count), x, y, Inches(2.4), Inches(0.6),
                      size=22, bold=True, colour=colour, family="Cambria")
            _add_text(slide, label, x, y + Inches(0.55), Inches(2.4), Inches(0.4),
                      size=10, colour=MUTED, family="Cascadia Code")

        _add_text(slide, f"{a.row_count} line item{'' if a.row_count == 1 else 's'}",
                  Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
                  size=12, colour=MUTED, italic=True)


def _add_outro(prs, result) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, INK)
    _add_text(slide, "End of pack", Inches(0.6), Inches(2.8),
              Inches(12), Inches(1.5), size=44, bold=True, colour=PAPER, family="Cambria")
    _add_text(
        slide,
        f"Generated by the Day 03 Budget Workbench for {result.metadata.company}.",
        Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
        size=14, colour=(0xC0, 0xB8, 0xA0),
    )


# ---- Building blocks --------------------------------------------------

def _add_text(slide, text, left, top, width, height, *,
              size=14, bold=False, italic=False, colour=INK,
              family="Segoe UI Variable Display", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = family
    run.font.color.rgb = RGBColor(*colour)
    return box


def _set_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colour)


def _gbp(v: float) -> str:
    if v is None:
        return "n/a"
    return f"GBP {v:,.0f}"


def _gbp_signed(v: float) -> str:
    if v is None:
        return "n/a"
    sign = "-" if v < 0 else ("+" if v > 0 else "")
    return f"{sign}GBP {abs(v):,.0f}"


def _pct_signed(v: float | None) -> str:
    if v is None:
        return "n/a"
    sign = "" if v < 0 else "+"
    return f"{sign}{v * 100:.1f}%"


def _bullets(items: list[str]) -> str:
    if not items:
        return "(none)"
    return "\n".join(f"• {item}" for item in items)


def _slug(s: str) -> str:
    s = s.strip().lower()
    s = re.sub(r"[^a-z0-9]+", "-", s).strip("-")
    return s or "untitled"
